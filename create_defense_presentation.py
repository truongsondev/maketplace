from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT_DIR = Path(r"D:\maketplace")
ASSET_DIR = OUT_DIR / "ppt_assets"
PPTX_PATH = OUT_DIR / "KLTN_Aura_Fashion_BaoVe.pptx"
SCRIPT_PATH = OUT_DIR / "KLTN_Aura_Fashion_BaoVe_script.md"


W, H = 13.333, 7.5
COLORS = {
    "ink": RGBColor(23, 28, 38),
    "muted": RGBColor(89, 99, 117),
    "paper": RGBColor(248, 249, 251),
    "line": RGBColor(218, 223, 232),
    "navy": RGBColor(22, 42, 73),
    "blue": RGBColor(29, 111, 216),
    "teal": RGBColor(25, 151, 146),
    "green": RGBColor(51, 142, 80),
    "orange": RGBColor(219, 126, 44),
    "red": RGBColor(202, 69, 69),
    "white": RGBColor(255, 255, 255),
    "soft_blue": RGBColor(232, 242, 255),
    "soft_teal": RGBColor(229, 247, 246),
    "soft_green": RGBColor(232, 246, 237),
    "soft_orange": RGBColor(255, 241, 226),
}


def add_textbox(slide, text, x, y, w, h, size=18, bold=False, color="ink",
                align=PP_ALIGN.LEFT, font="Aptos", margin=0.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return box


def add_title(slide, title, subtitle=None, kicker=None):
    if kicker:
        add_textbox(slide, kicker.upper(), 0.65, 0.32, 4.8, 0.28, 8.5, True, "blue")
    add_textbox(slide, title, 0.62, 0.62, 7.9, 0.78, 22, True, "ink")
    if subtitle:
        add_textbox(slide, subtitle, 0.65, 1.32, 8.5, 0.42, 11.5, False, "muted")
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(1.86), Inches(1.15), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["blue"]
    line.line.fill.background()


def add_footer(slide, n):
    add_textbox(slide, "KLTN - Aura Fashion Marketplace", 0.62, 7.08, 3.4, 0.22, 8.5, False, "muted")
    add_textbox(slide, f"{n:02d}", 12.1, 7.05, 0.5, 0.25, 8.5, True, "muted", PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title, body, accent="blue"):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1)
    add_textbox(slide, title, x + 0.18, y + 0.16, w - 0.36, 0.28, 11, True, accent)
    add_textbox(slide, body, x + 0.18, y + 0.52, w - 0.36, h - 0.62, 10.3, False, "ink")
    return shape


def pill(slide, text, x, y, w, color="blue", fill="soft_blue"):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.fill.background()
    add_textbox(slide, text, x + 0.08, y + 0.08, w - 0.16, 0.2, 8.7, True, color, PP_ALIGN.CENTER)


def bullet_list(slide, items, x, y, w, h, size=15, color="ink"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.04)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS[color]
        p.space_after = Pt(8)
    return box


def add_picture(slide, file_name, x, y, w=None, h=None, crop=False):
    path = ASSET_DIR / file_name
    if not path.exists():
        return None
    pic = slide.shapes.add_picture(
        str(path), Inches(x), Inches(y), width=Inches(w) if w else None, height=Inches(h) if h else None
    )
    return pic


def add_process(slide, steps, x, y, w, h):
    gap = 0.12
    sw = (w - gap * (len(steps) - 1)) / len(steps)
    accents = ["blue", "teal", "green", "orange", "blue"]
    fills = ["soft_blue", "soft_teal", "soft_green", "soft_orange", "soft_blue"]
    for i, (title, desc) in enumerate(steps):
        sx = x + i * (sw + gap)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(sx), Inches(y), Inches(sw), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS[fills[i % len(fills)]]
        shape.line.fill.background()
        add_textbox(slide, f"{i+1}", sx + 0.08, y + 0.12, 0.35, 0.3, 12, True, accents[i % len(accents)])
        add_textbox(slide, title, sx + 0.42, y + 0.12, sw - 0.52, 0.28, 9.6, True, "ink")
        add_textbox(slide, desc, sx + 0.18, y + 0.56, sw - 0.36, h - 0.66, 8.6, False, "muted")
        if i < len(steps) - 1:
            add_textbox(slide, ">", sx + sw + 0.02, y + h / 2 - 0.12, 0.08, 0.22, 12, True, "muted")


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["paper"]


def make_deck():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    slides_notes = []

    def new_slide(title, subtitle=None, kicker=None):
        slide = prs.slides.add_slide(blank)
        bg(slide)
        add_title(slide, title, subtitle, kicker)
        add_footer(slide, len(prs.slides))
        return slide

    # 1
    s = prs.slides.add_slide(blank)
    bg(s)
    add_picture(s, "image51.png", 7.15, 0.48, w=5.55)
    add_textbox(s, "KHÓA LUẬN TỐT NGHIỆP", 0.68, 0.58, 4.2, 0.28, 9, True, "blue")
    add_textbox(s, "Xây dựng website thời trang tích hợp AI phân tích hành vi khách hàng",
                0.65, 1.05, 6.2, 1.15, 26, True, "ink")
    add_textbox(s, "Aura Fashion Marketplace", 0.68, 2.45, 4.7, 0.38, 15, True, "teal")
    bullet_list(s, [
        "SVTH: Lê Trường Sơn - 22110407",
        "SVTH: Trần Tuấn - 22110448",
        "GVHD: PGS.TS. Hoàng Văn Dũng",
    ], 0.72, 3.15, 5.6, 1.1, 12.5, "ink")
    pill(s, "E-commerce thời trang", 0.72, 4.68, 1.8, "blue", "soft_blue")
    pill(s, "Recommendation", 2.7, 4.68, 1.5, "teal", "soft_teal")
    pill(s, "Chatbot", 4.35, 4.68, 1.0, "green", "soft_green")
    pill(s, "IDM-VTON", 5.5, 4.68, 1.15, "orange", "soft_orange")
    add_footer(s, 1)
    slides_notes.append(("1. Mở đầu", "Giới thiệu tên đề tài, nhóm thực hiện, định hướng chính: xây dựng website thời trang có các chức năng thương mại điện tử và tích hợp AI để cá nhân hóa trải nghiệm mua sắm."))

    # 2
    s = new_slide("Mục tiêu buổi trình bày", "Đi theo đúng các tiêu chí chấm: thực tiễn, phương pháp, sản phẩm và khả năng trả lời phản biện.", "Tổng quan")
    add_process(s, [
        ("Vấn đề", "Vì sao đề tài có tính thực tiễn trong bán hàng thời trang"),
        ("Giải pháp", "Kiến trúc website, nghiệp vụ, dữ liệu và luồng xử lý"),
        ("AI", "Recommendation, chatbot tư vấn và thử đồ IDM-VTON"),
        ("Kết quả", "Sản phẩm, kiểm thử, ưu điểm, hạn chế"),
        ("Phản biện", "Lập luận cho các lựa chọn kỹ thuật quan trọng"),
    ], 0.78, 2.35, 11.85, 2.1)
    slides_notes.append(("2. Mục tiêu", "Nói ngắn rằng bài trình bày sẽ không chỉ mô tả giao diện, mà tập trung vào bài toán, cách thiết kế, lý do chọn công nghệ và kết quả đạt được."))

    # 3
    s = new_slide("Bối cảnh và vấn đề", "Người mua thời trang cần chọn đúng sản phẩm, đúng phong cách và đúng ngữ cảnh sử dụng.", "Tính thực tiễn")
    card(s, 0.75, 2.15, 3.75, 2.55, "Khó khăn của khách hàng", "Nhiều mẫu mã, màu sắc, kích cỡ và phong cách khiến người dùng mất thời gian khi lựa chọn.", "blue")
    card(s, 4.78, 2.15, 3.75, 2.55, "Khó khăn của cửa hàng", "Cần quản lý sản phẩm, biến thể, đơn hàng, voucher, hoàn tiền, tồn kho và dashboard vận hành.", "teal")
    card(s, 8.82, 2.15, 3.75, 2.55, "Cơ hội ứng dụng AI", "AI có thể hỗ trợ tư vấn, gợi ý sản phẩm và thử đồ trực tuyến để tăng trải nghiệm mua sắm.", "orange")
    add_textbox(s, "Luận điểm chính: đề tài không chỉ là một website bán hàng, mà là nền tảng thử nghiệm cá nhân hóa trong thương mại điện tử thời trang.", 1.0, 5.35, 11.2, 0.55, 15, True, "ink", PP_ALIGN.CENTER)
    slides_notes.append(("3. Bối cảnh", "Nhấn mạnh tính thực tiễn: thời trang có nhiều biến thể và yếu tố cảm tính, nên hệ thống cần hỗ trợ quyết định mua hàng tốt hơn."))

    # 4
    s = new_slide("Mục tiêu và phạm vi đề tài", "Xây dựng hệ thống hoàn chỉnh từ khách hàng, quản trị đến AI service.", "Mục tiêu")
    bullet_list(s, [
        "Xây dựng website mua sắm thời trang cho khách hàng.",
        "Xây dựng trang quản trị cho admin/seller.",
        "Thiết kế backend API xử lý nghiệp vụ thương mại điện tử.",
        "Tích hợp PayOS, voucher, đánh giá, yêu thích, hoàn tiền.",
        "Tích hợp AI: gợi ý sản phẩm, chatbot tư vấn, thử đồ trực tuyến.",
    ], 0.95, 2.08, 5.7, 3.6, 15, "ink")
    add_picture(s, "image53.png", 7.0, 2.0, w=5.2)
    card(s, 7.0, 5.1, 5.2, 0.9, "Phạm vi", "Tập trung vào quy trình mua hàng, quản lý vận hành và cá nhân hóa trải nghiệm thời trang.", "green")
    slides_notes.append(("4. Mục tiêu", "Nêu rõ phạm vi vừa đủ: không phải sàn thương mại điện tử tổng hợp, mà là marketplace thời trang có định hướng mở rộng."))

    # 5
    s = new_slide("Chức năng chính", "Tách rõ nhóm người dùng khách hàng và nhóm quản trị.", "Sản phẩm")
    add_picture(s, "image55.png", 7.05, 1.98, w=5.3)
    card(s, 0.75, 2.0, 2.85, 2.35, "Khách hàng", "Đăng ký, đăng nhập, tìm kiếm, xem chi tiết, giỏ hàng, đặt hàng, thanh toán, đánh giá, yêu thích.", "blue")
    card(s, 3.8, 2.0, 2.85, 2.35, "Admin/Seller", "Quản lý sản phẩm, biến thể, đơn hàng, người dùng, voucher, banner, hoàn tiền và thống kê.", "teal")
    card(s, 0.75, 4.65, 2.85, 1.45, "AI", "Gợi ý sản phẩm, chatbot tư vấn, thử đồ bằng IDM-VTON.", "orange")
    card(s, 3.8, 4.65, 2.85, 1.45, "Vận hành", "Docker, Redis, RabbitMQ, Prometheus/Grafana, Nginx.", "green")
    slides_notes.append(("5. Chức năng", "Trình bày theo người dùng thật: khách mua hàng thao tác gì, admin/seller quản trị gì, AI bổ trợ ở đâu."))

    # 6
    s = new_slide("Kiến trúc tổng thể", "Monorepo gồm frontend khách hàng, frontend quản trị, backend API, AI service và hạ tầng.", "Thiết kế")
    card(s, 0.75, 2.05, 2.25, 1.15, "Client khách hàng", "Next.js, React, Tailwind", "blue")
    card(s, 0.75, 3.55, 2.25, 1.15, "Client quản trị", "Vite, React, Tailwind", "teal")
    card(s, 3.55, 2.75, 2.45, 1.45, "Backend API", "Node.js, Express, TypeScript, Prisma", "green")
    card(s, 6.55, 2.0, 2.25, 1.15, "MySQL", "Dữ liệu nghiệp vụ", "blue")
    card(s, 6.55, 3.35, 2.25, 1.15, "Redis/RabbitMQ", "Cache và xử lý bất đồng bộ", "orange")
    card(s, 9.35, 2.65, 2.55, 1.55, "AI Service", "FastAPI, pgvector, embedding, recommendation", "teal")
    add_textbox(s, "Frontend", 1.2, 5.3, 1.0, 0.25, 9, True, "muted", PP_ALIGN.CENTER)
    add_textbox(s, "API nghiệp vụ", 4.0, 5.3, 1.4, 0.25, 9, True, "muted", PP_ALIGN.CENTER)
    add_textbox(s, "Dữ liệu & hạ tầng", 6.85, 5.3, 1.7, 0.25, 9, True, "muted", PP_ALIGN.CENTER)
    add_textbox(s, "AI", 10.1, 5.3, 0.8, 0.25, 9, True, "muted", PP_ALIGN.CENTER)
    slides_notes.append(("6. Kiến trúc", "Nhấn mạnh sự tách lớp giúp dễ bảo trì, dễ mở rộng và phù hợp với hệ thống có AI service chạy độc lập."))

    # 7
    s = new_slide("Thiết kế nghiệp vụ và dữ liệu", "Hệ thống được thiết kế quanh vòng đời sản phẩm, đơn hàng và tương tác người dùng.", "Phương pháp")
    add_picture(s, "image3.png", 0.78, 2.05, h=4.35)
    card(s, 6.45, 2.05, 2.7, 1.5, "Nghiệp vụ lõi", "Product, variant, cart, order, payment, return, refund.", "blue")
    card(s, 9.45, 2.05, 2.7, 1.5, "Dữ liệu AI", "Recommendation events, embeddings, similarities, cache.", "teal")
    card(s, 6.45, 3.95, 2.7, 1.5, "Bảo mật", "JWT, session Redis, phân quyền admin/seller/user.", "green")
    card(s, 9.45, 3.95, 2.7, 1.5, "Vận hành", "Webhook PayOS, worker đối soát, metrics và logging.", "orange")
    slides_notes.append(("7. Thiết kế", "Nói rằng nhóm có phân tích use case, ERD, class/sequence diagram và thiết kế dữ liệu riêng cho recommendation."))

    # 8
    s = new_slide("Vai trò của AI trong hệ thống", "AI được dùng để hỗ trợ quyết định mua hàng, không thay thế nghiệp vụ thương mại điện tử.", "AI")
    card(s, 0.8, 2.15, 3.55, 2.6, "Recommendation", "Gợi ý sản phẩm tương tự, phù hợp hành vi, phổ biến hoặc theo ngữ cảnh trang.", "blue")
    card(s, 4.9, 2.15, 3.55, 2.6, "Chatbot tư vấn", "Hỗ trợ hỏi đáp sản phẩm, phong cách, thông tin cửa hàng và quy trình mua hàng.", "teal")
    card(s, 9.0, 2.15, 3.55, 2.6, "IDM-VTON", "Thử đồ trực tuyến giúp người dùng hình dung sản phẩm trên cơ thể tốt hơn.", "orange")
    add_textbox(s, "Điểm chung: các chức năng AI đều phục vụ trải nghiệm người mua thời trang, nơi quyết định mua phụ thuộc nhiều vào sở thích và hình dung cá nhân.", 1.15, 5.38, 10.9, 0.62, 14.5, True, "ink", PP_ALIGN.CENTER)
    slides_notes.append(("8. Vai trò AI", "Nêu rõ AI nằm ở tầng hỗ trợ trải nghiệm: gợi ý, tư vấn, thử đồ. Không nói quá rằng AI quyết định tất cả."))

    # 9
    s = new_slide("Luồng gợi ý sản phẩm", "Hybrid recommendation: kết hợp nội dung sản phẩm, hành vi người dùng, độ phổ biến và ngữ cảnh.", "Recommendation")
    add_picture(s, "image4.png", 0.75, 2.0, w=5.85)
    add_process(s, [
        ("Track", "Ghi nhận xem, tìm kiếm, yêu thích, giỏ hàng, mua hàng"),
        ("Store", "Lưu event, sản phẩm, embedding và cache"),
        ("Score", "Tính điểm tương đồng, hành vi, phổ biến, ngữ cảnh"),
        ("Fallback", "Dùng danh mục hoặc sản phẩm phổ biến khi thiếu dữ liệu"),
    ], 6.95, 2.15, 5.35, 2.15)
    card(s, 7.05, 4.75, 5.15, 1.15, "Giá trị", "Người dùng thấy sản phẩm liên quan nhanh hơn; hệ thống có nền tảng cá nhân hóa khi dữ liệu tăng.", "green")
    slides_notes.append(("9. Recommendation", "Giải thích đơn giản: không chỉ AI embedding, mà có fallback và nhiều tín hiệu. Điều này làm chức năng phù hợp cả khi dữ liệu ban đầu chưa lớn."))

    # 10
    s = new_slide("Phản biện: Có cần recommendation khi sản phẩm chưa quá nhiều?", "Lập luận bảo vệ lựa chọn chức năng gợi ý sản phẩm.", "Câu hỏi hội đồng")
    card(s, 0.75, 2.0, 3.65, 2.65, "Thừa nhận", "Nếu số sản phẩm rất ít và không có biến thể, recommendation chưa tạo khác biệt lớn.", "orange")
    card(s, 4.75, 2.0, 3.65, 2.65, "Phản biện chính", "Thời trang không chỉ nhiều sản phẩm; còn nhiều kiểu dáng, màu sắc, size, phong cách và nhu cầu phối đồ.", "blue")
    card(s, 8.75, 2.0, 3.65, 2.65, "Giá trị kỹ thuật", "Thiết kế sớm luồng tracking, cache, embedding và scoring giúp hệ thống sẵn sàng mở rộng.", "teal")
    add_textbox(s, "Câu trả lời ngắn: Recommendation trong đề tài là nền tảng cá nhân hóa cho website thời trang, không chỉ là công cụ lọc sản phẩm của sàn lớn.", 1.0, 5.25, 11.25, 0.6, 15, True, "ink", PP_ALIGN.CENTER)
    slides_notes.append(("10. Phản biện recommendation", "Trả lời theo cấu trúc: thừa nhận điều kiện đúng của thầy, sau đó nói đặc thù thời trang và giá trị nền tảng kỹ thuật."))

    # 11
    s = new_slide("Chatbot tư vấn thời trang", "Chatbot giúp người dùng hỏi đáp nhanh trong quá trình mua hàng.", "AI")
    bullet_list(s, [
        "Hỗ trợ tư vấn sản phẩm theo nhu cầu, phong cách và ngữ cảnh.",
        "Giảm thao tác tìm kiếm thủ công khi người dùng chưa biết chọn gì.",
        "Có thể kết hợp dữ liệu cửa hàng, thông tin sản phẩm và chính sách mua hàng.",
        "Là bước đầu để phát triển trợ lý mua sắm cá nhân hóa.",
    ], 0.95, 2.05, 5.7, 3.3, 15, "ink")
    add_picture(s, "image62.png", 7.0, 2.0, w=5.2)
    slides_notes.append(("11. Chatbot", "Nêu chatbot là hỗ trợ trải nghiệm, nhất là trong thời trang khi người dùng có câu hỏi về phối đồ, chất liệu, dịp sử dụng."))

    # 12
    s = new_slide("Thử đồ trực tuyến bằng mô hình IDM-VTON", "Giúp người dùng hình dung sản phẩm trước khi quyết định mua.", "AI")
    card(s, 0.75, 2.0, 3.55, 2.4, "Đầu vào", "Ảnh người dùng và ảnh sản phẩm thời trang.", "blue")
    card(s, 4.9, 2.0, 3.55, 2.4, "Xử lý", "Mô hình IDM-VTON sinh ảnh thử đồ ảo dựa trên người và trang phục.", "teal")
    card(s, 9.05, 2.0, 3.55, 2.4, "Đầu ra", "Ảnh minh họa sản phẩm khi mặc lên người dùng.", "orange")
    add_textbox(s, "Ý nghĩa với đề tài: giảm rào cản mua online, tăng mức độ tự tin khi chọn sản phẩm thời trang.", 1.2, 5.05, 10.8, 0.55, 15, True, "ink", PP_ALIGN.CENTER)
    slides_notes.append(("12. IDM-VTON", "Dùng tên tiêu đề học thuật: thử đồ trực tuyến bằng mô hình IDM-VTON. Nói rõ đây là chức năng hỗ trợ hình dung, không cam kết thay thế thử đồ thực tế."))

    # 13
    s = new_slide("Công nghệ triển khai", "Lựa chọn công nghệ theo từng tầng của hệ thống.", "Cài đặt")
    techs = [
        ("Frontend", "Next.js, React, TypeScript, Tailwind CSS", "blue", "soft_blue"),
        ("Admin", "Vite, React, TypeScript, Tailwind CSS", "teal", "soft_teal"),
        ("Backend", "Node.js, Express, Prisma, MySQL", "green", "soft_green"),
        ("Queue/Cache", "RabbitMQ, Redis", "orange", "soft_orange"),
        ("AI", "FastAPI, pgvector, embedding", "blue", "soft_blue"),
        ("DevOps", "Docker, Nginx, Prometheus, Grafana", "teal", "soft_teal"),
    ]
    for i, (t, b, c, f) in enumerate(techs):
        x = 0.85 + (i % 3) * 4.05
        y = 2.05 + (i // 3) * 1.75
        shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.45), Inches(1.18))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS[f]
        shape.line.fill.background()
        add_textbox(s, t, x + 0.2, y + 0.15, 3.05, 0.25, 11.2, True, c)
        add_textbox(s, b, x + 0.2, y + 0.52, 3.05, 0.35, 10.5, False, "ink")
    slides_notes.append(("13. Công nghệ", "Không cần đọc từng công nghệ quá dài; nhấn mạnh chọn công nghệ theo vai trò và đã tổ chức thành các service độc lập."))

    # 14
    s = new_slide("Kiểm thử và đánh giá", "Kiểm thử theo luồng nghiệp vụ chính và các tình huống vận hành.", "Chất lượng sản phẩm")
    add_picture(s, "image59.png", 7.0, 2.0, w=5.15)
    bullet_list(s, [
        "Kiểm thử đăng ký, đăng nhập, phân quyền và session.",
        "Kiểm thử luồng sản phẩm, giỏ hàng, đặt hàng, thanh toán.",
        "Kiểm thử admin/seller: sản phẩm, đơn hàng, voucher, dashboard.",
        "Kiểm thử recommendation: có dữ liệu, thiếu dữ liệu và fallback.",
        "Theo dõi metrics backend/AI service bằng Prometheus/Grafana.",
    ], 0.95, 2.05, 5.7, 3.5, 14.5, "ink")
    slides_notes.append(("14. Kiểm thử", "Nêu kiểm thử theo kịch bản người dùng và theo thành phần kỹ thuật. Đây là phần ăn điểm về mức độ hoàn thiện sản phẩm."))

    # 15
    s = new_slide("Kết quả đạt được", "Sản phẩm đã hoàn thiện các luồng chính và có nền tảng AI mở rộng.", "Kết quả")
    card(s, 0.75, 2.0, 3.7, 2.25, "Website khách hàng", "Mua hàng, thanh toán, theo dõi đơn, yêu thích, đánh giá và gợi ý sản phẩm.", "blue")
    card(s, 4.8, 2.0, 3.7, 2.25, "Trang quản trị", "Quản lý sản phẩm, đơn hàng, voucher, banner, người dùng, hoàn tiền, dashboard.", "teal")
    card(s, 8.85, 2.0, 3.7, 2.25, "AI service", "Recommendation, chatbot, thử đồ trực tuyến và pipeline dữ liệu hành vi.", "orange")
    card(s, 2.7, 4.75, 3.7, 1.2, "Hạ tầng", "Docker Compose, Redis, RabbitMQ, Nginx, Prometheus/Grafana.", "green")
    card(s, 6.95, 4.75, 3.7, 1.2, "Tài liệu", "Use case, sequence, database schema, báo cáo và tài liệu API.", "blue")
    slides_notes.append(("15. Kết quả", "Tóm kết những gì đã làm được. Tránh nói chung chung; gắn với sản phẩm thực tế và tài liệu kỹ thuật."))

    # 16
    s = new_slide("Hạn chế và hướng phát triển", "Nhìn nhận đúng giới hạn hiện tại và hướng nâng cấp sau bảo vệ.", "Đánh giá")
    card(s, 0.75, 2.0, 5.45, 2.95, "Hạn chế", "Dữ liệu hành vi thực tế còn ít nên cá nhân hóa chưa thể đánh giá bằng CTR/conversion. Chatbot cần thêm dữ liệu thời trang chi tiết. IDM-VTON phụ thuộc chất lượng ảnh đầu vào.", "orange")
    card(s, 6.85, 2.0, 5.45, 2.95, "Hướng phát triển", "Bổ sung A/B testing, dashboard KPI cho AI, tối ưu trọng số recommendation, mở rộng dữ liệu phong cách/chất liệu/form dáng và hoàn thiện quy trình production.", "green")
    add_textbox(s, "Tinh thần bảo vệ: nhóm không khẳng định hệ thống đã tối ưu tuyệt đối; nhóm chứng minh được kiến trúc, sản phẩm và nền tảng để tiếp tục phát triển.", 1.0, 5.55, 11.25, 0.62, 14.5, True, "ink", PP_ALIGN.CENTER)
    slides_notes.append(("16. Hạn chế", "Phần này nên nói chủ động, không né hạn chế. Hội đồng thường đánh giá cao khi nhóm biết điểm yếu và có hướng đo lường cải thiện."))

    # 17
    s = new_slide("Định vị theo tiêu chí chấm", "Các điểm cần nhấn mạnh khi trình bày và trả lời câu hỏi.", "Kết luận")
    card(s, 0.75, 2.0, 2.65, 2.4, "Thực tiễn", "Bài toán mua sắm thời trang online, quản trị vận hành và cá nhân hóa trải nghiệm.", "blue")
    card(s, 3.65, 2.0, 2.65, 2.4, "Phương pháp", "Khảo sát, phân tích yêu cầu, use case, thiết kế dữ liệu, kiến trúc service.", "teal")
    card(s, 6.55, 2.0, 2.65, 2.4, "Sản phẩm", "Website khách hàng, admin/seller, backend API, AI service, dashboard.", "green")
    card(s, 9.45, 2.0, 2.65, 2.4, "Phản biện", "Biết giải thích lý do chọn AI/recommendation và giới hạn hiện tại.", "orange")
    add_textbox(s, "Xin cảm ơn hội đồng đã lắng nghe", 2.1, 5.3, 9.2, 0.55, 24, True, "ink", PP_ALIGN.CENTER)
    slides_notes.append(("17. Kết thúc", "Kết bằng lời cảm ơn. Nếu demo trực tiếp, chuyển từ slide này sang demo hệ thống theo kịch bản đã chuẩn bị."))

    prs.save(PPTX_PATH)

    lines = ["# Kịch bản thuyết trình bảo vệ KLTN", ""]
    for title, note in slides_notes:
        lines.append(f"## {title}")
        lines.append(note)
        lines.append("")
    lines.append("## Gợi ý trả lời nhanh về recommendation")
    lines.append("Dạ em đồng ý nếu hệ thống chỉ có rất ít sản phẩm và không có biến thể thì recommendation chưa tạo khác biệt lớn. Tuy nhiên đề tài của nhóm là website thời trang có định hướng mở rộng. Trong thời trang, khó khăn không chỉ nằm ở số lượng sản phẩm mà còn ở màu sắc, size, kiểu dáng, phong cách và khả năng phối đồ. Vì vậy recommendation đóng vai trò cá nhân hóa trải nghiệm và tạo nền tảng tracking, cache, embedding, scoring ngay từ đầu. Khi dữ liệu tăng, hệ thống có thể tối ưu theo CTR, conversion rate và doanh thu từ gợi ý.")
    lines.append("")
    SCRIPT_PATH.write_text("\n".join(lines), encoding="utf-8")


if __name__ == "__main__":
    make_deck()
    print(PPTX_PATH)
    print(SCRIPT_PATH)
