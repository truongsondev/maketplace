from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(r"D:\maketplace")
OUT = ROOT / "AI_Recommendation_Aura_Deep_Dive.pptx"
SCRIPT = ROOT / "AI_Recommendation_Aura_Deep_Dive_Speaker_Notes.md"
W, H = 13.333, 7.5

C = {
    "bg": RGBColor(247, 249, 252), "ink": RGBColor(18, 28, 45),
    "muted": RGBColor(82, 96, 117), "line": RGBColor(214, 222, 233),
    "white": RGBColor(255, 255, 255), "blue": RGBColor(32, 103, 226),
    "cyan": RGBColor(17, 153, 173), "green": RGBColor(35, 151, 93),
    "orange": RGBColor(232, 131, 42), "red": RGBColor(211, 69, 91),
    "purple": RGBColor(113, 75, 190), "softblue": RGBColor(232, 241, 255),
    "softcyan": RGBColor(226, 247, 249), "softgreen": RGBColor(231, 247, 238),
    "softorange": RGBColor(255, 242, 224), "softred": RGBColor(255, 233, 238),
    "softpurple": RGBColor(241, 235, 255), "dark": RGBColor(12, 31, 58),
}

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
blank = prs.slide_layouts[6]
notes = []

def tb(slide, text, x, y, w, h, size=14, bold=False, color="ink", align=PP_ALIGN.LEFT, margin=.05, font="Aptos"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin); tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = C[color]
    return box

def rect(slide, x, y, w, h, fill="white", line="line", radius=True):
    typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(typ, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = C[fill]
    if line: s.line.color.rgb = C[line]
    else: s.line.fill.background()
    return s

def line(slide, x1, y1, x2, y2, color="line", width=1.5, arrow=False):
    s = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    s.line.color.rgb = C[color]; s.line.width = Pt(width)
    if arrow: s.line.end_arrowhead = True
    return s

def title(slide, text, sub="", kicker="AI RECOMMENDATION • TECHNICAL DEEP DIVE"):
    tb(slide, kicker, .65, .22, 6.8, .25, 8.5, True, "blue")
    tb(slide, text, .62, .55, 11.9, .55, 22, True)
    if sub: tb(slide, sub, .65, 1.08, 11.7, .4, 10.5, False, "muted")
    rect(slide, .65, 1.52, 1.15, .045, "blue", None, False)

def footer(slide):
    n = len(prs.slides)
    tb(slide, "Aura Fashion Marketplace • Recommendation System", .65, 7.15, 5, .18, 8, False, "muted")
    tb(slide, f"{n:02d}", 12.1, 7.12, .5, .2, 8.5, True, "muted", PP_ALIGN.RIGHT)

def new(text, sub="", note=""):
    s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = C["bg"]
    title(s, text, sub); footer(s); notes.append((text, note)); return s

def card(slide, x, y, w, h, head, body, accent="blue", fill="white", fs=10.5):
    rect(slide, x, y, w, h, fill, "line")
    rect(slide, x, y, .07, h, accent, None, False)
    tb(slide, head, x+.2, y+.15, w-.35, .28, 11, True, accent)
    tb(slide, body, x+.2, y+.52, w-.35, h-.62, fs, False, "ink")

def bullets(slide, items, x, y, w, h, size=12, color="ink"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(.08); tf.margin_right = Inches(.03); tf.margin_top = Inches(.03)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.text = "• " + item
        p.font.name = "Aptos"; p.font.size = Pt(size); p.font.color.rgb = C[color]; p.space_after = Pt(7)
    return box

def pill(slide, text, x, y, w, fill="softblue", color="blue"):
    rect(slide, x, y, w, .34, fill, None)
    tb(slide, text, x+.04, y+.075, w-.08, .18, 8.5, True, color, PP_ALIGN.CENTER)

def flow(slide, steps, y=2.2, x=.75, total=11.85, h=1.25):
    gap=.18; sw=(total-gap*(len(steps)-1))/len(steps)
    fills=["softblue","softcyan","softgreen","softorange","softpurple","softred"]
    accents=["blue","cyan","green","orange","purple","red"]
    for i,(a,b) in enumerate(steps):
        xx=x+i*(sw+gap); rect(slide,xx,y,sw,h,fills[i%6],None)
        tb(slide,str(i+1),xx+.1,y+.12,.28,.25,11,True,accents[i%6])
        tb(slide,a,xx+.42,y+.12,sw-.52,.28,10,True)
        tb(slide,b,xx+.14,y+.53,sw-.28,h-.58,8.7,False,"muted")
        if i<len(steps)-1: tb(slide,"→",xx+sw+.02,y+.47,.14,.25,13,True,"muted",PP_ALIGN.CENTER)

def table(slide, headers, rows, x, y, widths, row_h=.42, fs=9.2, highlight=None):
    xx=x
    for h,w in zip(headers,widths):
        rect(slide,xx,y,w,row_h,"dark",None,False); tb(slide,h,xx+.06,y+.1,w-.12,.2,8.8,True,"white",PP_ALIGN.CENTER); xx+=w
    for ri,row in enumerate(rows):
        xx=x; fill="softgreen" if highlight is not None and ri==highlight else ("white" if ri%2==0 else "softblue")
        for val,w in zip(row,widths):
            rect(slide,xx,y+(ri+1)*row_h,w,row_h,fill,"line",False); tb(slide,str(val),xx+.06,y+(ri+1)*row_h+.095,w-.12,.22,fs,ri==highlight,"ink",PP_ALIGN.CENTER); xx+=w

# 1 cover
s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=C["dark"]
tb(s,"AURA FASHION",.72,.55,3.2,.3,9,True,"cyan")
tb(s,"AI Recommendation\nTừ tín hiệu hành vi đến một sản phẩm được xếp hạng #1",.7,1.22,7.1,1.5,28,True,"white")
tb(s,"Deep dive vào candidate generation, embedding, pgvector, hybrid scoring, merge, exclusion, cache và feedback loop trong chính dự án.",.75,3.0,6.4,.85,13,False,"white")
for i,(s1,s2,c,f) in enumerate([("384D","Embedding","blue","softblue"),("0.75 / 0.25","Hybrid","cyan","softcyan"),("4 feeds","Serving","green","softgreen")]):
    xx=8.2+i*1.55
    rect(s,xx,1.3,1.35,1.25,f,None); tb(s,s1,xx+.08,1.57,1.19,.3,13.5,True,c,PP_ALIGN.CENTER); tb(s,s2,xx+.08,2.0,1.19,.2,9,True,"ink",PP_ALIGN.CENTER)
tb(s,"Nội dung bám sát: Node.js recommendation repository + FastAPI AI service + MySQL/Redis/RabbitMQ + PostgreSQL/pgvector",.75,5.65,11.4,.55,10.5,False,"white")
tb(s,"2026 • Technical presentation",.75,6.65,4,.25,9,True,"cyan")
notes.append(("AI Recommendation", "Mở đầu: mục tiêu của bài không phải kể API nào trả danh sách, mà giải phẫu vì sao một candidate cụ thể nhận điểm, vượt qua bộ lọc và đứng đầu."))

# 2
s=new("Câu hỏi trung tâm: vì sao sản phẩm này được gợi ý?","Không bắt đầu từ endpoint; bắt đầu từ bằng chứng và phép tính.")
card(s,.7,1.85,3.8,3.9,"1 • Bằng chứng người dùng","Người dùng xem váy ren, thêm váy chữ A vào giỏ, tìm “váy trắng dự tiệc”. Mỗi hành vi có trọng số và cửa sổ thời gian khác nhau.","blue","white",13)
card(s,4.75,1.85,3.8,3.9,"2 • Bằng chứng sản phẩm","Tên, mô tả, danh mục và tags được biến thành vector 384 chiều. Khoảng cách cosine đo mức gần về ngữ nghĩa.","cyan","white",13)
card(s,8.8,1.85,3.8,3.9,"3 • Quyết định xếp hạng","Nhiều nguồn sinh candidate; hệ thống loại sản phẩm đã mua/đang ở giỏ, hợp nhất điểm, sắp xếp và cache kết quả.","green","white",13)
pill(s,"Output = candidate hợp lệ có bằng chứng mạnh nhất tại placement hiện tại",3.8,6.15,5.75,"softpurple","purple")

# 3
s=new("Kiến trúc triển khai thực tế","Hai lớp ranking: orchestration trong Node.js và semantic retrieval trong AI service.")
flow(s,[("Storefront","phát event + yêu cầu feed"),("RabbitMQ","buffer/retry/DLQ"),("Node repository","candidate + rules + cache"),("FastAPI AI","embedding + hybrid"),("pgvector","ANN cosine search"),("Redis/MySQL","signals + artifacts")],2.0,.55,12.2,1.35)
card(s,.75,3.75,3.75,2.25,"Online path","Redis cache → nếu miss: SQL/Redis candidates → AI hybrid → merge → product cards → cache.","blue")
card(s,4.8,3.75,3.75,2.25,"Near-real-time path","Event vào RabbitMQ; consumer ghi MySQL, tăng hot score và invalidates cache liên quan.","orange")
card(s,8.85,3.75,3.75,2.25,"Offline/periodic path","Mỗi ~15 phút: co-occurrence từ order_items, train popularity/preferences, embed tối đa 2.000 sản phẩm.","purple")

# 4
s=new("Dữ liệu đầu vào: một event không chỉ là một click","Schema RecommendationEvent giữ identity, context, placement và chống trùng.")
table(s,["Trường","Ý nghĩa kỹ thuật","Ví dụ"],[
    ("eventType","tín hiệu implicit/explicit","VIEW_PRODUCT"),("userId","cá nhân hóa đăng nhập","user-123"),("sessionId","cold-start/anonymous context","sess-a9"),("productId","đối tượng hành vi","dress-rumy"),("searchQuery","ý định bằng ngôn ngữ","váy trắng"),("placement","đo hiệu quả vị trí","home_recommendations"),("dedupeKey","idempotency unique","view:...:timestamp"),("occurredAt","recency window","2026-07-12")],.8,1.85,[1.55,5.4,4.55],.47,10)

# 5
s=new("Ingestion đáng tin cậy: event đi qua RabbitMQ","Tách latency của UI khỏi persistence và feature update.")
flow(s,[("Publish","topic exchange"),("Consume","prefetch = 20"),("Persist","unique dedupeKey"),("Feature update","Redis ZSET/LIST"),("Invalidate","xóa feed cache")],1.9,.9,11.55,1.35)
card(s,.8,3.65,3.7,2.25,"Retry policy","Retry queue TTL 15 giây; tối đa 3 lần; sau đó dead-letter queue. Message durable và deliveryMode=2.","orange")
card(s,4.8,3.65,3.7,2.25,"Idempotency","Duplicate DB key trả accepted=true, duplicated=true: retry không làm hot score/persistence nhân đôi sau insert thành công.","green")
card(s,8.8,3.65,3.7,2.25,"Realtime freshness","Event mới tăng hot score, lưu 30 sản phẩm gần nhất/user và xóa cache home/cart/personalized liên quan.","blue")

# 6
s=new("Trọng số hành vi: giá trị kinh doanh được mã hóa thành số","Các pipeline dùng trọng số hơi khác nhau theo mục tiêu feed.")
table(s,["Event","Hot Redis","Trending 30d","Personalized 45d","AI train"],[
    ("PURCHASE","5","5","4 × 1.35","4"),("ADD_TO_CART","3","3","3 × 1.35","3"),("FAVORITE","2","2","2.5 × 1.35","2"),("VIEW","1","1","1 × 1.35","1"),("REC_CLICK","1.5","—","0.5 × 1.35","1"),("IMPRESSION","0","—","0.5 × 1.35","1")],1.0,2.0,[2.2,2.0,2.2,2.5,2.0],.55,10)
tb(s,"Ý nghĩa",1.0,5.95,1.1,.25,11,True,"blue"); tb(s,"Purchase mạnh hơn view; personalized còn boost 35% tín hiệu thuộc user. Đây là heuristic, chưa phải trọng số được học từ dữ liệu.",2.0,5.88,9.9,.55,11,False,"ink")

# 7
s=new("Offline refresh: biến lịch sử thành artifacts phục vụ nhanh","Scheduler chạy ngay lúc bootstrap và lặp theo RECOMMENDATION_REFRESH_INTERVAL_MS (mặc định 15 phút).")
flow(s,[("MySQL catalog","≤ 2.000 products"),("Recent events","≤ 5.000 events"),("Co-occurrence","self-join order_items"),("AI /train","preferences + popularity"),("/embed/products","upsert pgvector")],1.9,.7,11.9,1.4)
card(s,.8,3.75,5.7,2.2,"Co-occurrence artifact","score(A,B) = số đơn chứa đồng thời A và B. Lưu product_similarities với algorithm='cooccurrence', score, rank, metadata refreshedAt.","purple",fs=12)
card(s,6.8,3.75,5.7,2.2,"Semantic artifact","document = title + description + category + tags/attributes → normalized embedding 384D → PostgreSQL vector index.","cyan",fs=12)

# 8
s=new("Tạo document cho một sản phẩm","Embedding tốt hay xấu phụ thuộc trực tiếp vào chuỗi đầu vào.")
rect(s,.8,1.85,11.8,1.25,"dark",None); tb(s,'"Váy ren dáng A Rumy" + "cổ bèo, dự tiệc..." + "Váy nữ" + "color trắng occasion party material lace"',1.05,2.15,11.3,.55,14,True,"white",PP_ALIGN.CENTER)
flow(s,[("Normalize text","ghép field"),("Tokenizer","MiniLM hoặc whitespace"),("Encoder","384 numbers"),("L2 normalize","||v|| = 1"),("Persist","vector + metadata")],3.55,.9,11.55,1.25)
card(s,1.0,5.2,5.3,1.05,"Điểm mạnh","Tên + mô tả + category + tags giúp semantic retrieval hiểu “váy tiệc” gần “đầm dự sự kiện”.","green",fs=10)
card(s,7.0,5.2,5.3,1.05,"Giới hạn hiện tại","Variant attributes chi tiết chưa được sync đầy đủ; payload Node chủ yếu đưa tags vào attributes.","red",fs=10)

# 9
s=new("Encoder: MiniLM nếu tải được, hashing nếu không","Hệ thống có graceful fallback nhưng chất lượng hai chế độ rất khác nhau.")
card(s,.8,1.9,5.65,3.6,"all-MiniLM-L6-v2","SentenceTransformer encode(text, normalize_embeddings=True). Học semantic từ corpus lớn; từ đồng nghĩa có thể gần nhau dù token khác nhau. Output 384 chiều.","blue","softblue",13)
card(s,6.85,1.9,5.65,3.6,"hashing-baseline","Tách theo whitespace → hash(token) mod 384 → cộng tần suất → L2 normalize. Nhanh và offline, nhưng collision và không hiểu đồng nghĩa/ngữ cảnh.","orange","softorange",13)
tb(s,"Rủi ro vận hành: health.model_loaded=false không làm service chết, nhưng ranking semantic có thể giảm mạnh mà người dùng chỉ thấy chất lượng kém.",1.0,5.95,11.2,.55,11,True,"red",PP_ALIGN.CENTER)

# 10
s=new("Cosine similarity trong pgvector","Tìm vector gần context nhất bằng toán tử <=>.")
rect(s,.85,1.85,5.6,3.75,"white","line"); tb(s,"cos(q,p) = (q · p) / (||q|| ||p||)",1.2,2.35,4.9,.5,20,True,"blue",PP_ALIGN.CENTER)
tb(s,"Vì vector đã normalize L2:\ncos(q,p) ≈ q · p\n\nSQL score = 1 − (embedding <=> query_vector)",1.15,3.25,5,.95,14,False,"ink",PP_ALIGN.CENTER)
card(s,6.85,1.85,5.55,1.05,"0.92","Rất gần về nội dung/ngữ nghĩa","green","softgreen",11)
card(s,6.85,3.12,5.55,1.05,"0.65","Có một phần thuộc tính/danh mục chung","orange","softorange",11)
card(s,6.85,4.39,5.55,1.05,"0.10","Hầu như không liên quan","red","softred",11)
pill(s,"IVFFlat + vector_cosine_ops: giảm phạm vi tìm kiếm so với full scan",3.4,6.15,6.5,"softpurple","purple")

# 11
s=new("Context vector: hệ thống hiểu 'gu hiện tại' như thế nào?","Không embed user trực tiếp trong request; lấy trung bình vector các sản phẩm context.")
flow(s,[("Recent events","top 10 / 45 ngày"),("Context IDs","view/cart/fav/buy"),("Product vectors","v₁...vₙ"),("Mean","q = mean(vᵢ)"),("Normalize","q / ||q||")],1.9,.85,11.65,1.35)
rect(s,1.1,3.75,11.05,1.05,"dark",None); tb(s,"q_user = normalize((v_váy-ren + v_váy-chữ-A + v_áo-cổ-bèo) / 3)",1.3,4.08,10.65,.35,17,True,"white",PP_ALIGN.CENTER)
bullets(s,["Nếu user đăng nhập: AI còn nối tối đa 8 sản phẩm preference phổ biến vào context.","user_profile tuổi/chiều cao/cân nặng được gửi sang AI nhưng hiện chưa tham gia công thức recommend().","Nếu không có context vector: AI vector_hits rỗng, chỉ còn popularity fallback."],1.15,5.15,10.9,1.2,10.5)

# 12
s=new("Công thức hybrid trong AI service","Semantic relevance chiếm 75%; demand evidence chiếm 25%.")
rect(s,.8,1.8,11.7,1.05,"dark",None); tb(s,"final_AI(p) = 0.75 × cosine(q,p) + 0.25 × log(1 + popularity(p))",1.1,2.12,11.1,.36,18,True,"white",PP_ALIGN.CENTER)
table(s,["Candidate","Cosine","Raw popularity","log1p(pop)","Final AI"],[
    ("Đầm dự tiệc B","0.91","4","1.609","1.085"),("Váy basic C","0.72","20","3.045","1.301"),("Áo ren D","0.83","1","0.693","0.796")],1.0,3.35,[3.0,1.7,2.0,2.0,2.0],.58,10,1)
tb(s,"Kết luận ví dụ",1.0,5.95,1.6,.25,11,True,"orange"); tb(s,"C vượt B dù semantic thấp hơn, vì popularity chưa được normalize và log1p có thể > 1. Đây là chủ đề calibration quan trọng.",2.5,5.86,9.6,.65,11,True,"ink")

# 13
s=new("Candidate generation: không có một 'model duy nhất'","Mỗi nguồn trả productId + score + reason + source.")
cards=[("pgvector","semantic/content","blue"),("co-occurrence","mua cùng đơn","purple"),("hot ZSET","real-time demand","red"),("trending SQL","30-day behavior","orange"),("search intent","query match","cyan"),("category/type","business fallback","green"),("recent behavior","personal context","blue"),("latest/sale","cold-start fallback","orange")]
for i,(a,b,c) in enumerate(cards):
    x=.75+(i%4)*3.08; y=1.85+(i//4)*1.6; card(s,x,y,2.8,1.3,a,b,c,fs=10)
rect(s,2.2,5.45,8.9,.7,"softpurple",None); tb(s,"Candidate generation ưu tiên recall: lấy limit × 2/3/4 rồi mới merge, exclude và cắt top-K.",2.45,5.68,8.4,.28,12,True,"purple",PP_ALIGN.CENTER)

# 14
s=new("Feed Home: xu hướng + phiên anonymous + ý định tìm kiếm","Mục tiêu là relevance tức thời ngay cả khi chưa đăng nhập.")
flow(s,[("Hot Redis","realtime weighted"),("Trending SQL","30 ngày"),("Search intent","30 ngày/session"),("Session behavior","14 ngày"),("Merge max+20%","top limit")],1.9,.8,11.75,1.35)
table(s,["Session event","Weight"],[('ADD_TO_CART','3'),('FAVORITE_PRODUCT','2.5'),('VIEW_PRODUCT','1'),('khác','0.5')],1.05,3.65,[3.3,2.0],.48,10)
card(s,6.8,3.65,5.4,2.25,"Cache & fallback","Cache theo session + limit, TTL 900s. Nếu không có context: sản phẩm có ảnh, ưu tiên isSale rồi createdAt mới nhất.","blue",fs=12)

# 15
s=new("Feed Product Detail: 'sản phẩm tương tự' được tạo thế nào?","Context là đúng một product đang xem; output luôn loại chính product đó.")
flow(s,[("Current product","context ID"),("AI vector","limit × 2"),("Co-occurrence","limit × 2"),("Category/type","limit × 2"),("Merge + exclude","top K")],1.9,.85,11.65,1.35)
card(s,.9,3.65,3.55,2.2,"Semantic","Hai sản phẩm có mô tả/thuộc tính gần nhau dù chưa từng được mua cùng.","cyan",fs=11.5)
card(s,4.85,3.65,3.55,2.2,"Behavioral","Hai sản phẩm xuất hiện trong cùng đơn: bằng chứng người mua thực tế.","purple",fs=11.5)
card(s,8.8,3.65,3.55,2.2,"Fallback","Cùng category, parent/children hoặc productType; ưu tiên sale và mới cập nhật.","green",fs=11.5)

# 16
s=new("Feed Cart: cross-sell từ nhiều sản phẩm context","Không gợi ý lại bất kỳ product nào đang nằm trong giỏ.")
flow(s,[("Cart items","excluded + context"),("AI mean vector","limit × 3"),("Related pairs","co-occurrence"),("Category family","parent/child"),("uniqueTop","top K")],1.9,.85,11.65,1.35)
rect(s,.95,3.7,11.4,1.15,"dark",None); tb(s,"Ví dụ: giỏ = {váy ren, túi trắng} → context mean mang cả occasion + style → candidate giày cao gót có thể gần cả hai.",1.25,4.0,10.8,.55,14,True,"white",PP_ALIGN.CENTER)
bullets(s,["AI reason: “phù hợp với các sản phẩm trong giỏ”.","Co-occurrence cung cấp cross-sell dựa trên đơn thật.","TTL 600 giây; ADD/REMOVE event invalidates cart cache."],1.15,5.25,10.8,1.0,11)

# 17
s=new("Feed Personalized: pipeline đầy đủ nhất","Kết hợp behavior 45 ngày, search intent, AI, co-occurrence và category.")
flow(s,[("Score behavior","PURCHASE 4..."),("Exclude","context+cart+90d bought"),("AI retrieval","limit × 4"),("Other sources","search/related/category"),("Fallback","latest if short")],1.8,.7,11.95,1.35)
card(s,.8,3.5,5.6,2.35,"Own signals","recent event score × 1.2 được đưa lại làm candidate, nhưng contextIds đồng thời nằm trong exclusion nên thực tế các own candidates này thường bị loại ở uniqueTop.","red",fs=11.3)
card(s,6.8,3.5,5.6,2.35,"Exclusion","Không gợi lại sản phẩm context, sản phẩm đang trong cart và mọi sản phẩm đã mua trong 90 ngày—giảm lặp nhưng có thể bỏ lỡ replenishment.","green",fs=11.3)

# 18
s=new("Search intent: chuyển câu người dùng thành candidate","Hiện dùng lexical SQL LIKE, chưa embed query.")
flow(s,[("Collect query","30 ngày"),("Rank terms","count + recency"),("LIKE match","name/description"),("Rank penalty","term & product index"),("uniqueTop","candidate list")],1.9,.85,11.65,1.35)
rect(s,.9,3.65,11.55,1.1,"softcyan",None); tb(s,"score = max(0.5, query_count × 3 − termIndex × 0.75 − productIndex × 0.25)",1.15,3.98,11.05,.35,16,True,"cyan",PP_ALIGN.CENTER)
card(s,1.0,5.15,5.3,1.0,"Tốt","Phản ứng nhanh với “váy trắng”, dễ giải thích.","green",fs=10)
card(s,7.0,5.15,5.3,1.0,"Chưa tốt","Không hiểu synonym/typo; query score có thang khác cosine.","red",fs=10)

# 19
s=new("Co-occurrence: học quan hệ từ giỏ hàng đã mua","Một dạng item-item collaborative signal đơn giản, minh bạch.")
rect(s,.8,1.85,5.5,3.9,"white","line"); tb(s,"SQL self-join",1.05,2.1,2,.28,12,True,"purple")
tb(s,"order_items oi1\nJOIN order_items oi2\n  ON same order_id\n AND product_id khác nhau\n\nscore(A,B) = COUNT(*)",1.1,2.65,4.9,1.8,15,True,"ink",PP_ALIGN.CENTER,font="Courier New")
card(s,6.7,1.85,5.65,1.05,"A + B cùng 12 đơn","score(A,B)=12","purple","softpurple",11)
card(s,6.7,3.15,5.65,1.05,"A + C cùng 4 đơn","score(A,C)=4","blue","softblue",11)
card(s,6.7,4.45,5.65,1.05,"Vấn đề popularity bias","Không normalize theo độ phổ biến A/B; bestseller dễ liên quan với mọi thứ.","red","softred",10.5)

# 20
s=new("Merge scores trong Node: max thắng, nguồn phụ cộng 20%","Đây là bước cuối quyết định candidate nào đứng trên candidate nào.")
rect(s,.8,1.8,11.7,.95,"dark",None); tb(s,"Nếu product xuất hiện nhiều lần: giữ source có score lớn nhất; mỗi score nhỏ hơn đến sau cộng thêm 0.2 × score đó.",1.1,2.07,11.1,.4,14,True,"white",PP_ALIGN.CENTER)
table(s,["Product","pgvector","cooccur","category","Kết quả merge"],[
    ("B","1.085","0.80","—","1.085 + .2×.80 = 1.245"),("C","1.301","—","8.0","8.0 + .2×1.301 = 8.260"),("D","0.796","1.20","7.0","7.0 + .2×1.2 + .2×.796 = 7.399")],.9,3.2,[2.0,1.8,1.8,1.8,4.0],.58,9.5,1)
tb(s,"Technical debt",1.0,5.75,1.5,.25,11,True,"red"); tb(s,"Score cosine (~0–1), co-occurrence (count), fallback rank (limit…1), search (count×3) không cùng thang; category fallback có thể lấn át AI.",2.35,5.65,9.8,.65,11,True,"ink")

# 21
s=new("Bộ lọc và guardrail trước khi hiển thị","Ranking tốt nhưng candidate không hợp lệ vẫn phải bị loại.")
cards=[("Context","không gợi lại item đang xem","blue"),("Cart","không gợi item trong giỏ","cyan"),("Purchased 90d","giảm lặp personalized","purple"),("Soft delete","product isDeleted=false","red"),("Card lookup","missing product bị drop","orange"),("Unique","mỗi product một lần","green")]
for i,(a,b,c) in enumerate(cards): card(s,.9+(i%3)*4.05,1.9+(i//3)*1.75,3.7,1.4,a,b,c,fs=11)
rect(s,2.1,5.55,9.1,.65,"softred",None); tb(s,"Khoảng trống hiện tại: getProductCards chưa lọc Product.status ACTIVE, variant stock/status hay stockAvailable > 0.",2.35,5.78,8.6,.25,11,True,"red",PP_ALIGN.CENTER)

# 22
s=new("Worked example #1: AI chọn candidate nào?","Context: user vừa xem váy ren trắng và váy chữ A; candidate là B/C/D.")
table(s,["Candidate","Cosine","Events pop","log1p","0.75×cos","0.25×log","AI score"],[
    ("B Đầm dự tiệc","0.91","4","1.609","0.683","0.402","1.085"),("C Váy basic","0.72","20","3.045","0.540","0.761","1.301"),("D Áo ren","0.83","1","0.693","0.623","0.173","0.796")],.55,2.0,[2.5,1.35,1.35,1.35,1.5,1.5,1.5],.62,9.3,1)
card(s,.8,4.85,3.65,1.1,"#1 C • 1.301","Demand thắng semantic","orange","softorange",10)
card(s,4.85,4.85,3.65,1.1,"#2 B • 1.085","Semantic cao nhất","blue","softblue",10)
card(s,8.9,4.85,3.65,1.1,"#3 D • 0.796","Ít popularity","cyan","softcyan",10)
tb(s,"Nếu business muốn semantic chi phối thật sự, popularity cần normalize (min-max/percentile) trước khi nhân 0.25.",1.0,6.25,11.2,.35,11,True,"purple",PP_ALIGN.CENTER)

# 23
s=new("Worked example #2: sau merge, thứ hạng có thể đảo","AI score không phải final score của Node repository.")
table(s,["Product","Nguồn mạnh nhất","Nguồn phụ","Final Node","Rank"],[
    ("C Váy basic","category=8.0","AI=1.301","8.260","#1"),("D Áo ren","category=7.0","co=1.2; AI=.796","7.399","#2"),("B Đầm tiệc","AI=1.085","co=.80","1.245","#3")],1.0,2.0,[2.5,3.0,3.1,1.8,1.2],.66,10,0)
flow(s,[("Semantic winner","B: cosine .91"),("AI winner","C: score 1.301"),("Node winner","C: category 8"),("Displayed #1","C after exclude")],4.8,1.25,10.8,1.1)
tb(s,"Thông điệp bảo vệ",2.4,6.25,2,.25,11,True,"blue"); tb(s,"“AI retrieval tạo candidate tốt; orchestration + heuristic calibration mới quyết định sản phẩm cuối cùng.”",4.2,6.15,6.8,.5,12,True,"ink",PP_ALIGN.CENTER)

# 24
s=new("End-to-end: một sản phẩm đi từ catalog đến vị trí #1","Chuỗi quyết định đầy đủ, không bỏ qua tầng dữ liệu hay serving.")
flow(s,[("Catalog text","title/desc/category/tags"),("Embedding","MiniLM 384D"),("User context","mean recent vectors"),("Retrieve","cosine top 3K"),("Hybrid",".75 semantic+.25 pop"),("Merge/filter","multi-source top K")],1.8,.55,12.2,1.35)
cards=[("Bước 1","C được cosine=.72","cyan"),("Bước 2","pop=20 → log=3.045","orange"),("Bước 3","AI=1.301","blue"),("Bước 4","category rank=8","green"),("Bước 5","merge=8.260","purple"),("Bước 6","không bị exclude → #1","red")]
for i,(a,b,c) in enumerate(cards): card(s,.7+i*2.08,3.75,1.85,1.45,a,b,c,fs=9.5)
rect(s,2.0,5.75,9.3,.62,"dark",None); tb(s,"Sản phẩm được gợi ý không vì “API trả về”, mà vì sống sót qua 6 phép biến đổi và guardrail.",2.25,5.96,8.8,.25,12,True,"white",PP_ALIGN.CENTER)

# 25
s=new("Cache hai tầng: tốc độ mà vẫn có khả năng phục hồi","Redis là fast path; MySQL RecommendationCache là durable fallback.")
flow(s,[("Request","cache key by context"),("Redis HIT","return immediately"),("Redis MISS","DB cache lookup"),("DB HIT","warm Redis"),("Full MISS","generate + persist")],1.9,.8,11.75,1.35)
table(s,["Feed","Cache key","TTL"],[
    ("Home","home:v2:{session}:{limit}","900s"),("Product","product:{product}:{limit}","1800s"),("Cart","cart:{user}:{limit}","600s"),("Personalized","personalized:v4:{user}:{session}:{limit}","900s")],1.25,3.75,[2.2,6.8,1.6],.48,10)

# 26
s=new("Cache invalidation: freshness đến từ event","Không chờ TTL nếu hành vi vừa thay đổi.")
table(s,["Event context","Patterns bị xóa"],[
    ("sessionId","home:{session}:*, home:v2:{session}:*, personalized:*:{session}:*"),("userId","cart:{user}:*, personalized v1/v2/v3/v4 của user")],1.1,2.0,[2.5,8.7],.7,10)
card(s,1.0,4.0,5.3,1.65,"Ưu điểm","ADD_TO_CART làm feed cart/personalized thay đổi ngay ở request kế tiếp.","green",fs=12)
card(s,7.0,4.0,5.3,1.65,"Rủi ro","Dùng Redis KEYS(pattern) có thể block khi keyspace lớn; production nên chuyển SCAN hoặc versioned keys.","red",fs=12)

# 27
s=new("Fault tolerance: AI hỏng thì feed không nhất thiết hỏng","Node repository bắt exception và tiếp tục bằng co-occurrence/category/catalog.")
flow(s,[("AI request","/recommend/hybrid"),("Timeout/5xx","throw"),("Catch","increment fallback metric"),("Return []","no AI candidates"),("Fallback sources","still produce feed")],1.9,.85,11.65,1.35)
card(s,.9,3.7,3.55,2.0,"Resilience","AI service failure không làm endpoint recommendation 500 nếu các nguồn khác còn dữ liệu.","green",fs=11.5)
card(s,4.85,3.7,3.55,2.0,"Consistency","RabbitMQ retry/DLQ giữ event ingestion; DB cache cứu khi Redis mất cache.","blue",fs=11.5)
card(s,8.8,3.7,3.55,2.0,"Blind spot","fetch AI chưa có explicit timeout/AbortController; request có thể treo lâu tùy network stack.","red",fs=11.5)

# 28
s=new("Observability: biết hệ thống nhanh hay chỉ 'có chạy'","Prometheus metrics được đo ở event, cache, feed và AI.")
table(s,["Metric","Labels","Câu hỏi trả lời"],[
    ("recommendation_events_total","event_type, source","Có nhận đúng hành vi?"),("cache_hits_total","feed","Cache hiệu quả bao nhiêu?"),("generation_latency_ms","feed","P95 feed nào chậm?"),("ai_latency_ms","operation","Train/embed/recommend chậm?"),("ai_fallback_total","feed","AI hỏng hoặc suy giảm?")],.75,1.9,[3.2,3.0,5.0],.58,9.5)
pill(s,"Nên bổ sung: candidate count/source, zero-result rate, score distribution, CTR/CVR theo model version",2.25,5.75,8.9,"softorange","orange")

# 29
s=new("Đánh giá chất lượng: offline metric chưa đủ","Recommendation phải đo cả relevance, diversity và tác động kinh doanh.")
card(s,.8,1.85,3.7,3.9,"Offline ranking","Recall@K: item mua có nằm trong candidate?\nNDCG@K: item đúng đứng cao?\nMRR: vị trí item đầu tiên đúng.\nCoverage: % catalog được expose.","blue",fs=12)
card(s,4.8,1.85,3.7,3.9,"Online behavior","Impression → click CTR.\nClick → add-to-cart.\nRecommendation-assisted purchase.\nRevenue/session và conversion uplift.","green",fs=12)
card(s,8.8,1.85,3.7,3.9,"Guardrail","Latency P95/P99.\nOut-of-stock exposure.\nDiversity category/brand.\nPopularity concentration và novelty.","orange",fs=12)

# 30
s=new("A/B testing đúng cách","Schema RecommendationExperiment đã có, nhưng serving chưa gắn assignment/model variant.")
flow(s,[("Stable hash","user/session → bucket"),("Assign","control vs treatment"),("Log impression","experiment+variant"),("Track outcome","click/cart/buy"),("Compare","uplift + significance")],1.9,.85,11.65,1.35)
card(s,.9,3.75,5.5,2.05,"Control","Heuristic hiện tại: unnormalized merge + 0.75/0.25 AI.","blue",fs=12)
card(s,6.9,3.75,5.5,2.05,"Treatment đề xuất","Normalize per-source → weighted RRF/learning-to-rank; diversity reranker; stock guardrail.","purple",fs=12)

# 31
s=new("Technical debt nhìn thẳng vào code","Những điểm cần nói rõ để bài bảo vệ có chiều sâu và trung thực.")
issues=[("Score calibration","nguồn dùng thang điểm khác nhau","red"),("Body profile unused","được gửi nhưng AI không dùng","orange"),("User embeddings unused","schema có, serving chưa đọc","purple"),("Product status/stock","card filter chưa đủ guardrail","red"),("Hash reproducibility","Python hash đổi theo process seed","orange"),("IVFFlat tuning","chưa thấy ANALYZE/lists/probes","purple"),("Own signal conflict","recent items bị exclusion loại","blue"),("No query embedding","search chỉ LIKE","cyan")]
for i,(a,b,c) in enumerate(issues): card(s,.7+(i%4)*3.1,1.8+(i//4)*1.72,2.85,1.42,a,b,c,fs=9.6)
rect(s,2.2,5.55,8.9,.67,"softgreen",None); tb(s,"Điểm mạnh: hệ thống có đường tiến hóa rõ—không cần phá kiến trúc để thay heuristic bằng model tốt hơn.",2.45,5.78,8.4,.25,11,True,"green",PP_ALIGN.CENTER)

# 32
s=new("Roadmap kỹ thuật: từ heuristic hybrid đến ranking học được","Ưu tiên theo impact/risk, không chạy theo độ phức tạp mô hình.")
flow(s,[("P0 Guardrail","ACTIVE + stock > 0"),("P1 Normalize","per-source percentile"),("P2 RRF","rank fusion ổn định"),("P3 Diversity","MMR/category cap"),("P4 LTR","learn weights from logs")],1.85,.75,11.85,1.35)
table(s,["Giai đoạn","Thay đổi","Kết quả mong đợi"],[
    ("1–2 tuần","metrics + guardrail + SCAN cache","an toàn production"),("3–4 tuần","normalize + calibrated weights","ranking ổn định"),("5–8 tuần","A/B + MMR + query embedding","CTR/diversity tăng"),(">8 tuần","LightGBM/XGBoost LTR","trọng số học từ outcome")],1.0,3.65,[2.0,5.0,4.2],.52,9.8)

# 33
s=new("Kịch bản demo kỹ thuật","Chứng minh từng tầng bằng dữ liệu thay vì chỉ nhìn carousel.")
steps=[("1","Track VIEW/ADD","kiểm tra recommendation_events + Redis hot score"),("2","Refresh artifacts","kiểm tra co-occurrence + product embeddings"),("3","Query pgvector","so cosine của 3 candidate"),("4","Call feed","đọc score/reason/source"),("5","Trigger cache","request lần 2 và metric hit"),("6","New event","cache invalidated, ranking đổi")]
for i,(n,a,b) in enumerate(steps):
    y=1.8+i*.76; pill(s,n,.8,y,.42,"softblue","blue"); tb(s,a,1.4,y+.04,2.2,.25,10.5,True,"ink"); tb(s,b,3.65,y+.04,8.5,.28,10.5,False,"muted")
rect(s,2.4,6.45,8.5,.42,"dark",None); tb(s,"Demo tốt nhất kết thúc bằng việc giải thích vì sao rank thay đổi sau một ADD_TO_CART.",2.65,6.56,8,.18,9.5,True,"white",PP_ALIGN.CENTER)

# 34
s=new("Kết luận: recommendation là một chuỗi quyết định có thể kiểm chứng","Từ product text và hành vi đến candidate cuối cùng đều có dữ liệu, công thức và fallback.")
card(s,.8,1.9,3.7,3.65,"Retrieval","Embedding 384D + cosine/pgvector mở rộng semantic recall; co-occurrence và category cứu cold-start/AI failure.","cyan","softcyan",13)
card(s,4.8,1.9,3.7,3.65,"Ranking","Hybrid 0.75/0.25 tạo AI score; Node merge max+20%, exclusion và fallback quyết định final rank.","purple","softpurple",13)
card(s,8.8,1.9,3.7,3.65,"Learning loop","Event → RabbitMQ → MySQL/Redis → refresh artifacts → cache invalidation → recommendation mới.","green","softgreen",13)
tb(s,"Một sản phẩm đứng #1 vì nó có bằng chứng phù hợp nhất sau toàn bộ pipeline—not because an API happened to return it.",1.1,6.15,11.1,.45,14,True,"blue",PP_ALIGN.CENTER)

# speaker notes markdown and ppt notes
for slide, (head, note) in zip(prs.slides, notes):
    text = note or f"Trình bày slide “{head}”: liên hệ trực tiếp với code hiện tại và nhấn mạnh các con số/công thức trên slide."
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass

prs.save(OUT)
md=["# Speaker notes — AI Recommendation Aura Deep Dive", "", "Các ghi chú dưới đây bám theo thứ tự slide.", ""]
for i,(head,note) in enumerate(notes,1):
    md += [f"## Slide {i:02d} — {head}", "", note or "Giải thích sơ đồ/công thức trên slide, đối chiếu với pipeline thực tế và nêu trade-off kỹ thuật.", ""]
SCRIPT.write_text("\n".join(md), encoding="utf-8")
print(OUT)
print(SCRIPT)
print(f"slides={len(prs.slides)}")
